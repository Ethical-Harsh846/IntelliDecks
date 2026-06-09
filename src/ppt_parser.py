from pptx import Presentation


def extract_projects_from_pptx(
        pptx_path,
        slides_per_project=2,
        skip_first_slide=True):

    prs = Presentation(pptx_path)

    slides = list(prs.slides)

    # Skip cover slide
    if skip_first_slide:
        slides = slides[1:]

    all_projects = []

    for i in range(0, len(slides), slides_per_project):

        slide_group = slides[i:i + slides_per_project]

        slide_numbers = [
            x + 2 if skip_first_slide else x + 1
            for x in range(i, i + len(slide_group))
        ]

        project_id = ""
        project_title = ""
        domain = ""
        tech_stack = ""
        team_size = ""
        duration = ""
        project_type = ""

        all_text_parts = []

        # Extract text from both slides
        for slide in slide_group:

            for shape in slide.shapes:

                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()

                if text:
                    all_text_parts.append(text)

        if not all_text_parts:
            continue

        # First slide structure
        project_id = all_text_parts[0].strip()

        if len(all_text_parts) > 1:
            project_title = all_text_parts[1].strip()

        # Find PROJECT DETAILS section
        project_details_idx = -1

        for idx, line in enumerate(all_text_parts):
            if line.strip() == "PROJECT DETAILS":
                project_details_idx = idx
                break

        if project_details_idx != -1:

            search_window = all_text_parts[
                project_details_idx:
                min(project_details_idx + 15, len(all_text_parts))
            ]

            for idx, line in enumerate(search_window):

                line = line.strip()

                if line == "Domain" and idx + 1 < len(search_window):
                    domain = search_window[idx + 1].strip()

                elif line == "Team Size" and idx + 1 < len(search_window):
                    team_size = search_window[idx + 1].strip()

                elif line == "Duration" and idx + 1 < len(search_window):
                    duration = search_window[idx + 1].strip()

                elif line == "Type" and idx + 1 < len(search_window):
                    project_type = search_window[idx + 1].strip()

                elif line == "TECHNOLOGY STACK" and idx + 1 < len(search_window):
                    tech_stack = search_window[idx + 1].strip()

        combined_text = "\n".join(all_text_parts)

        all_projects.append({
            "project_id": project_id,
            "project_title": project_title,
            "domain": domain,
            "team_size": team_size,
            "duration": duration,
            "project_type": project_type,
            "tech_stack": tech_stack,
            "slide_numbers": slide_numbers,
            "combined_text": combined_text
        })

    print(f" Extracted {len(all_projects)} projects")

    return all_projects